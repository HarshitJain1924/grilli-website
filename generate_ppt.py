from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define Colors
    ORANGE = RGBColor(255, 140, 0)
    DARK_GREY = RGBColor(40, 40, 40)
    WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_GREY)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Grilli Website"
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    title.text_frame.paragraphs[0].font.size = Pt(60)

    subtitle.text = "Next-Gen Enterprise Kubernetes Transformation\nHigh Availability | Auto-Scaling | Real-time Monitoring"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE

    # --- Slide 2: Modern Microservices Architecture ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    slide.shapes.title.text = "Core Architecture"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "The Grilli project is orchestrated into four mission-critical layers:"
    
    p = tf.add_paragraph()
    p.text = "• Frontend: Responsive Nginx-based UI served via LoadBalancer."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Backend: High-performance Node.js API with HPA support."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Database: Persistent MongoDB with Dedicated Volume Claims (PVC)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Infrastructure: Metrics Server for real-time cluster intelligence."
    p.level = 1

    # --- Slide 3: High Availability (HA) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Engineered for Reliability"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Eliminating Single Points of Failure:"
    
    p = tf.add_paragraph()
    p.text = "• 2x Multi-Pod Redundancy: Minimum 2 replicas per service at all times."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Self-Healing: Liveness & Readiness probes automatically restart unhealthy containers."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Zero-Downtime Updates: Rolling update strategy for seamless deployments."
    p.level = 1

    # --- Slide 4: Dynamic Auto-Scaling (HPA) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Intelligent Scaling"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Adapting to customer demand in real-time:"
    
    p = tf.add_paragraph()
    p.text = "• Backend Scaling: Automatically expands to 10 replicas if CPU hits 70%."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Frontend Scaling: Expands to 5 replicas if CPU usage reaches 80%."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Smart Resource Management: Guaranteed CPU/Memory requests for every pod."
    p.level = 1

    # --- Slide 5: Observability & Monitoring ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Real-time Visibility"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "The Professional Observability Stack:"
    
    p = tf.add_paragraph()
    p.text = "• Prometheus: Scraping deep metrics from the Node.js runtime and K8s API."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Grafana Dashboards: Dark-themed visualization for traffic, errors, and memory."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Isolated Namespace: Monitoring stack runs independently in 'monitoring' namespace."
    p.level = 1

    # --- Slide 6: The Kubernetes Advantage ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, ORANGE)
    
    title = slide.shapes.title
    title.text = "Ready for the Future"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    left = top = Inches(3)
    width = height = Inches(4)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = "Grilli is now fully automated, scalable, and resilient."
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Save the presentation
    output_path = "Grilli_Kubernetes_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
